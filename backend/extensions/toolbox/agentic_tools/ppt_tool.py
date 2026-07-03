# Copyright (c) 2025 Bytedance Ltd. and/or its affiliates
# SPDX-License-Identifier: MIT

import json
import logging
import os
import re
import uuid
from pathlib import Path
from typing import Literal, Optional

from pydantic import BaseModel, Field

from langchain_core.messages import HumanMessage, SystemMessage
from langchain_core.tools import StructuredTool

from extensions._core.llms.llm import get_llm_by_type

from extensions.toolbox.agentic_tools.image_gen_tool import generate_image_to_path, is_image_gen_configured

# slide_deck 流程（engine=slide_deck 时委托）
try:
    from extensions.toolbox.agentic_tools.slide_deck_tool import _run_slide_deck
except ImportError:
    _run_slide_deck = None

logger = logging.getLogger(__name__)

PPT_OUTPUT_DIR = "ppt_output"
OUTLINE_SYSTEM = """你是一个专业的PPT大纲撰写助手。根据用户给出的一句话主题，生成一份结构清晰的PPT大纲。
要求：
1. 使用 Markdown 格式输出。
2. 第一行用一级标题（# ）写整份PPT的标题。
3. 每一页幻灯片用一个二级标题（## ）作为该页标题，紧跟的列表项（- 或 *）为该页要点，每页3-6个要点为宜。
4. 总页数控制在5-10页（含封面）。
5. 只输出大纲内容，不要输出其他解释。"""


SLIDE_IMAGE_PROMPT_SYSTEM = """你是一个幻灯片配图助手。根据给定的幻灯片标题和要点，生成一句英文的图片描述（prompt），用于 DALL-E 等文生图模型。
要求：一句话，描述该页幻灯片主题对应的简洁、专业、适合做配图的画面，不要包含多余解释。只输出这一句英文。"""


class GeneratePPTInput(BaseModel):
    """Input for the PPT generation tool (outline or generate)."""

    engine: Literal["simple", "slide_deck"] = Field(
        default="slide_deck",
        description="'simple': 文字型 PPT；'slide_deck': baoyu 风格整页图片幻灯片（默认）。",
    )
    topic: str = Field(default="", description="One-sentence topic for the PPT; required when action is 'outline'.")
    outline: str = Field(default="", description="Outline content; required when action is 'generate' (from step 1 or edited by user).")
    action: Literal["outline", "generate"] = Field(
        default="outline",
        description="'outline' to generate outline from topic; 'generate' to create .pptx from outline. (Used when engine='simple'.)",
    )
    with_images: bool = Field(
        default=False,
        description="When engine='simple' and action='generate', whether to generate an image per slide (requires IMAGE_GEN).",
    )
    # slide_deck 专用参数（engine='slide_deck' 时生效）
    content: str = Field(default="", description="长文/Markdown 内容，与 topic 二选一；用于 slide_deck。")
    style: str = Field(default="blueprint", description="slide_deck 风格：blueprint, minimal, hand-drawn 等。")
    audience: str = Field(default="general", description="slide_deck 受众：general, beginners, experts, executives 等。")
    lang: str = Field(default="auto", description="slide_deck 语言：zh, en, auto。")
    slides: int = Field(default=8, ge=3, le=20, description="slide_deck 目标页数。")
    outline_only: bool = Field(default=False, description="slide_deck：仅生成并返回 outline。")
    prompts_only: bool = Field(default=False, description="slide_deck：生成 outline 与 prompts 后即返回。")
    images_only: bool = Field(default=False, description="slide_deck：生成图片后即返回，不合并 pptx/pdf。")


def _generate_outline(topic: str) -> str:
    """Use LLM to generate PPT outline from one-sentence topic."""
    if not (topic and topic.strip()):
        return json.dumps({"error": "topic is required when action is 'outline'"}, ensure_ascii=False)
    try:
        llm = get_llm_by_type("basic")
        messages = [
            SystemMessage(content=OUTLINE_SYSTEM),
            HumanMessage(content=f"请根据以下主题生成PPT大纲：\n\n{topic.strip()}"),
        ]
        response = llm.invoke(messages)
        outline_text = response.content if hasattr(response, "content") else str(response)
        outline_text = (outline_text or "").strip()
        return json.dumps({"outline": outline_text}, ensure_ascii=False)
    except Exception as e:
        logger.exception("PPT outline generation failed")
        return json.dumps({"error": f"生成大纲失败: {e!s}"}, ensure_ascii=False)


def _parse_outline_to_slides(outline: str) -> list[dict]:
    """Parse markdown outline into list of slides: [{"title": str, "bullets": [str]}, ...]."""
    lines = [s.strip() for s in outline.strip().split("\n") if s.strip()]
    slides = []
    current_title = None
    current_bullets = []

    for line in lines:
        if line.startswith("# "):
            title = line.lstrip("# ").strip()
            if not title:
                continue
            if current_title is not None:
                slides.append({"title": current_title, "bullets": current_bullets})
            current_title = title
            current_bullets = []
        elif line.startswith("## "):
            title = line.lstrip("# ").strip()
            if not title:
                continue
            if current_title is not None:
                slides.append({"title": current_title, "bullets": current_bullets})
            current_title = title
            current_bullets = []
        elif line.startswith("- ") or line.startswith("* "):
            bullet = line[2:].strip()
            if bullet:
                current_bullets.append(bullet)
        else:
            if current_title and line and not line.startswith("#"):
                current_bullets.append(line)

    if current_title is not None:
        slides.append({"title": current_title, "bullets": current_bullets})

    if not slides:
        first_line = lines[0] if lines else ""
        first_line = re.sub(r"^#+\s*", "", first_line).strip()
        if first_line:
            slides = [{"title": first_line, "bullets": []}]
            for line in lines[1:]:
                line = line.strip()
                if line and (line.startswith("- ") or line.startswith("* ")):
                    slides[0]["bullets"].append(line[2:].strip())
                elif line and not line.startswith("#"):
                    slides.append({"title": line, "bullets": []})
    return slides


def _slide_image_prompt(title: str, bullets: list) -> str:
    """用 LLM 根据幻灯片标题和要点生成一句英文配图描述。"""
    llm = get_llm_by_type("basic")
    content = f"标题: {title}\n要点: " + "; ".join((bullets or [])[:5])
    messages = [
        SystemMessage(content=SLIDE_IMAGE_PROMPT_SYSTEM),
        HumanMessage(content=content),
    ]
    response = llm.invoke(messages)
    text = (response.content if hasattr(response, "content") else str(response)) or ""
    return text.strip()[:500]


def _generate_pptx_from_outline(outline: str, with_images: bool = False) -> str:
    """Create .pptx from outline and return JSON with download_url and filename. Optionally add generated image per slide."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return json.dumps({"error": "python-pptx 未安装，请运行: uv add python-pptx"}, ensure_ascii=False)

    if not (outline and outline.strip()):
        return json.dumps({"error": "outline is required when action is 'generate'"}, ensure_ascii=False)

    slides_data = _parse_outline_to_slides(outline)
    if not slides_data:
        return json.dumps({"error": "无法从大纲解析出幻灯片内容，请检查格式（# 标题，## 页标题，- 要点）"}, ensure_ascii=False)

    use_images = with_images and is_image_gen_configured()
    image_paths: list[Optional[str]] = []
    if use_images:
        for i, slide_data in enumerate(slides_data):
            title = (slide_data.get("title") or "").strip() or ""
            bullets = slide_data.get("bullets") or []
            path_rel: Optional[str] = None
            try:
                prompt = _slide_image_prompt(title, bullets)
                if prompt:
                    path_rel = generate_image_to_path(prompt, size="1024x1024")
            except Exception as e:
                logger.warning("幻灯片 %d 配图失败，跳过: %s", i + 1, e)
            image_paths.append(path_rel)
    else:
        image_paths = [None] * len(slides_data)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    cwd = Path(os.getcwd()).resolve()

    for i, slide_data in enumerate(slides_data):
        title = (slide_data.get("title") or "").strip() or f"Slide {i + 1}"
        bullets = slide_data.get("bullets") or []

        if i == 0 and len(slides_data) > 1 and not bullets:
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            if slide.placeholders[1]:
                slide.placeholders[1].text = ""
            if image_paths[i]:
                try:
                    full_path = cwd / image_paths[i]
                    if full_path.is_file():
                        slide.shapes.add_picture(str(full_path), Inches(5.5), Inches(2), Inches(4), Inches(3.75))
                except Exception as e:
                    logger.warning("标题页插入图片失败: %s", e)
        else:
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = title
            body = slide.placeholders[1].text_frame
            body.clear()
            for b in bullets[:8]:
                p = body.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(14)
            if image_paths[i]:
                try:
                    full_path = cwd / image_paths[i]
                    if full_path.is_file():
                        slide.shapes.add_picture(str(full_path), Inches(5.5), Inches(1.8), Inches(4), Inches(3.75))
                except Exception as e:
                    logger.warning("内容页 %d 插入图片失败: %s", i + 1, e)

    out_dir = cwd / PPT_OUTPUT_DIR
    out_dir.mkdir(parents=True, exist_ok=True)
    filename = f"ppt_{uuid.uuid4().hex[:12]}.pptx"
    filepath = out_dir / filename
    prs.save(str(filepath))
    relative_path = f"{PPT_OUTPUT_DIR}/{filename}"
    download_url = f"/api/workspace-file?path={relative_path}"
    return json.dumps(
        {"download_url": download_url, "filename": filename, "path": relative_path},
        ensure_ascii=False,
    )


def _run_generate_ppt(
    engine: Literal["simple", "slide_deck"] = "slide_deck",
    topic: str = "",
    outline: str = "",
    action: Literal["outline", "generate"] = "outline",
    with_images: bool = False,
    content: str = "",
    style: str = "blueprint",
    audience: str = "general",
    lang: str = "auto",
    slides: int = 8,
    outline_only: bool = False,
    prompts_only: bool = False,
    images_only: bool = False,
) -> str:
    """StructuredTool 会按 schema 字段以关键字参数调用，因此签名需与 GeneratePPTInput 一致。"""
    engine = (engine or "slide_deck").strip() or "slide_deck"
    if engine == "slide_deck" and _run_slide_deck is not None:
        return _run_slide_deck(
            topic=topic,
            content=content,
            outline=outline,
            style=style,
            audience=audience,
            lang=lang,
            slides=slides,
            outline_only=outline_only,
            prompts_only=prompts_only,
            images_only=images_only,
        )
    # engine == "simple"
    topic = (topic or "").strip()
    outline = (outline or "").strip()
    action = (action or "outline").strip() or "outline"
    if action == "outline":
        return _generate_outline(topic)
    if action == "generate":
        return _generate_pptx_from_outline(outline, with_images=with_images)
    return json.dumps({"error": f"invalid action: {action}"}, ensure_ascii=False)


generate_ppt_tool = StructuredTool(
    name="generate_ppt_tool",
    description="Generate a PPT: engine='slide_deck' (default) for baoyu-style full-page image slides (outline→prompts→images→pptx/pdf); engine='simple' for text PPT (action=outline with topic, then action=generate with outline, optional with_images).",
    args_schema=GeneratePPTInput,
    func=_run_generate_ppt,
)
