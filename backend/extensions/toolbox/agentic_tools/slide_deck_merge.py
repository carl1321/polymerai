# Copyright (c) 2025 Bytedance Ltd. and/or its affiliates
# SPDX-License-Identifier: MIT

"""
纯 Python 实现：将 slide_deck 目录下的 01-slide-*.png 按顺序合并为 PPTX 与 PDF。
一图一页，铺满幻灯片/页面，无 Node/bun 依赖。
"""

import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

# 幻灯片图片文件名模式：01-slide-xxx.png，按数字排序
SLIDE_IMAGE_PATTERN = re.compile(r"^(\d+)-slide-.+\.png$", re.IGNORECASE)


def _sorted_slide_images(dir_path: Path) -> list[Path]:
    """收集目录下 01-slide-*.png 并按编号排序。"""
    files: list[tuple[int, Path]] = []
    for f in dir_path.iterdir():
        if not f.is_file():
            continue
        m = SLIDE_IMAGE_PATTERN.match(f.name)
        if m:
            num = int(m.group(1))
            files.append((num, f))
    files.sort(key=lambda x: x[0])
    return [p for _, p in files]


def merge_images_to_pptx(slide_deck_dir: Path, output_stem: str) -> Path:
    """
    将 slide_deck_dir 下 01-slide-*.png 按顺序一图一页生成 PPTX，铺满幻灯片（16:9）。
    输出文件：slide_deck_dir / {output_stem}.pptx
    返回输出文件的 Path。
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
    except ImportError:
        raise RuntimeError("请安装 python-pptx: uv add python-pptx")

    images = _sorted_slide_images(slide_deck_dir)
    if not images:
        raise FileNotFoundError(f"目录下未找到 01-slide-*.png: {slide_deck_dir}")

    prs = Presentation()
    # 16:9 幻灯片尺寸（Office 默认 16:9 为 13.333 x 7.5 inch）
    prs.slide_width = Emu(12192000)  # 13.333 in
    prs.slide_height = Emu(6858000)  # 7.5 in
    blank_layout = prs.slide_layouts[6]  # 空白版式

    for img_path in images:
        slide = prs.slides.add_slide(blank_layout)
        # 图片铺满整页：左边 0，上边 0，宽高等于幻灯片宽高
        slide.shapes.add_picture(
            str(img_path),
            Emu(0),
            Emu(0),
            prs.slide_width,
            prs.slide_height,
        )

    out_path = slide_deck_dir / f"{output_stem}.pptx"
    prs.save(str(out_path))
    logger.info("已生成 PPTX: %s", out_path)
    return out_path


def merge_images_to_pdf(slide_deck_dir: Path, output_stem: str) -> Path:
    """
    将 slide_deck_dir 下 01-slide-*.png 按顺序一图一页合并为 PDF。
    输出文件：slide_deck_dir / {output_stem}.pdf
    返回输出文件的 Path。
    """
    try:
        import img2pdf
    except ImportError:
        raise RuntimeError("请安装 img2pdf: uv add img2pdf")

    images = _sorted_slide_images(slide_deck_dir)
    if not images:
        raise FileNotFoundError(f"目录下未找到 01-slide-*.png: {slide_deck_dir}")

    out_path = slide_deck_dir / f"{output_stem}.pdf"
    with open(out_path, "wb") as f:
        f.write(img2pdf.convert([str(p) for p in images]))
    logger.info("已生成 PDF: %s", out_path)
    return out_path
